import os

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from langchain_chroma import Chroma
from langchain.storage import InMemoryStore
from langchain.retrievers.multi_vector import MultiVectorRetriever
from langchain_core.documents import Document
from chromadb import HttpClient

from api.utils.logger import logger
from api.utils.id_generator import gen_document_id
from .image_services import gen_image_summaries
from api.chatbot.agents import TEXT_COLLECTION_NAME, SUMMARY_COLLECTION_NAME, get_vector_store

IMAGE_URL_PREFIX = "./data/ppt_images"


class PPTLoader:
    """PPT file loader, load, indexing and save it into vector store"""

    def __init__(self, store: InMemoryStore, vs_client: HttpClient):
        self.text_vector_store = get_vector_store(TEXT_COLLECTION_NAME)
        self.summary_vector_store = get_vector_store(SUMMARY_COLLECTION_NAME)
        self.ppt_url = None
        self.store = store

    def extract_shape_content(
        self,
        shape: BaseShape,
        slide_index: int,
        slide_texts: list[str],
        slide_images: list[str],
    ):
        if hasattr(shape, "text"):
            slide_texts.append(shape.text)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shape.shapes:
                self.extract_shape_content(
                    shape=s,
                    slide_index=slide_index,
                    lide_texts=slide_texts,
                    slide_images=slide_images,
                )
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = shape.image
            # ---get image "file" contents---
            image_bytes = image.blob
            # ---assign a name for the file, e.g. 'image.jpg'---
            ppt_file_name = self.ppt_url.split("/")[-1].split(".")[0].lower()
            image_index = len(slide_images) + 1
            image_filename = f"{IMAGE_URL_PREFIX}/{ppt_file_name}_slide{slide_index}_image{image_index}.{image.ext}"
            with open(image_filename, "wb") as f:
                f.write(image_bytes)
            slide_images.append(image_filename)
            logger.info(
                f"Extracted image {image_filename} from ppt file {self.ppt_url}"
            )

    def extract_contents(self, file_url: str) -> tuple[list[str], list[dict[str, str]]]:
        self.ppt_url = file_url
        ppt = Presentation(file_url)
        extracted_texts = []
        extracted_images = []

        if not os.path.exists(IMAGE_URL_PREFIX):
            os.makedirs(IMAGE_URL_PREFIX)

        try:
            for index, slide in enumerate(ppt.slides):
                slide_texts = []
                slide_images = []
                for shape in slide.shapes:
                    self.extract_shape_content(
                        shape=shape,
                        slide_index=index,
                        slide_texts=slide_texts,
                        slide_images=slide_images,
                    )
                # Combine texts in one slide into one text, as normally texts in one slide are revelent
                # later use this as chunk put into vector database
                texts_in_slides = " ".join(slide_texts)
                extracted_texts.append(texts_in_slides)
                for image_url in slide_images:
                    extracted_images.append({"image_url": image_url, "related_info": texts_in_slides})
        except Exception as e:
            logger.exception(
                f"Failed to extract content from ppt file {file_url}: {e}"
            )
            raise

        return extracted_texts, extracted_images

    def load(self, file_url: str) -> bool:
        try:
            extracted_texts, extracted_images = self.extract_contents(file_url)
            # Embedding all extracted texts, one record per slide
            documents = [
                Document(page_content=text, metadata={"source": "ppt"})
                for text in extracted_texts
            ]
            ids = [gen_document_id() for _ in range(len(documents))]
            self.text_vector_store.add_documents(documents=documents, ids=ids)

            # Create multi vector retriever to save image summary embeddings and raw image files
            image_base64_list, image_summaries = gen_image_summaries(
                extracted_images
            )
            id_key = "image_id"
            multi_retriever = MultiVectorRetriever(
                vectorstore=self.summary_vector_store,
                docstore=self.store,
                id_key=id_key,
            )
            image_ids = [gen_document_id() for _ in range(len(extracted_images))]
            summary_docs = [
                Document(page_content=summary, metadata={id_key: image_ids[index]})
                for index, summary in enumerate(image_summaries)
            ]
            multi_retriever.vectorstore.add_documents(summary_docs)
            multi_retriever.docstore.mset(list(zip(image_ids, image_base64_list)))
        except Exception as e:
            logger.exception(
                f"Failed to load ppt file {file_url}: {e}"
            )
            return False
        
        return True
